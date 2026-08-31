"""개요 도판의 편집용 PPTX — 보고서 그림 1과 같은 내용을 파워포인트로 낸다.

목적
    사용자가 문구·배치를 직접 손볼 수 있게 도형과 텍스트를 개별 개체로 만든다.
    미니 패널은 report_overview_figure.py 가 저장한 panel_{A..F}.png 를 그대로 쓴다.
    따라서 그림 1과 PPTX의 수치는 항상 같은 산출물에서 온다.

선행
    python scripts/4_visualization/report_overview_figure.py   (panel_*.png 생성)

산출
    outputs/report/report_overview_editable.pptx   (슬라이드 1장, 16:9)

실행
    python scripts/4_visualization/report_overview_pptx.py
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Pt

ROOT = Path(__file__).resolve().parents[2]
PANEL = ROOT / "outputs" / "figures" / "00_overview"
# 사용자가 손본 report_overview_editable.pptx 를 덮지 않도록 별도 파일로 낸다.
OUT = ROOT / "outputs" / "report" / "report_overview_editable_v2.pptx"

# ---------------------------------------------------------------- 색(보고서 규약)
NAVY = RGBColor(0x1D, 0x4A, 0x78)
MID = RGBColor(0x2F, 0x6F, 0x9F)
GREY_H = RGBColor(0x6F, 0x77, 0x80)
INK = RGBColor(0x1A, 0x1A, 0x1A)
SUB = RGBColor(0x33, 0x41, 0x4D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GOAL_BG = RGBColor(0xEE, 0xF3, 0xF7)
CARD_BG = [RGBColor(0xF4, 0xF5, 0xF6), RGBColor(0xEF, 0xF4, 0xF8), RGBColor(0xF0, 0xF5, 0xF9)]
CARD_EC = [RGBColor(0xCD, 0xD1, 0xD5), RGBColor(0xBF, 0xD0, 0xDE), RGBColor(0xC3, 0xD4, 0xE0)]
EDGE = RGBColor(0xC3, 0xCC, 0xD4)

FONT = "Pretendard"          # 없으면 파워포인트가 대체 폰트로 렌더
FONT_FALLBACK = "맑은 고딕"

# ---------------------------------------------------------------- 슬라이드 치수(cm)
SW, SH = 33.87, 19.05        # 16:9
M = 1.2
GOAL_H = 1.5
HEAD_H = 1.15
GAP = 0.5
COL_GAP = 0.55
COL_W = (SW - 2 * M - 2 * COL_GAP) / 3.0
PAD = 0.5
LAB_H = 0.85
CAP_H = 0.95
# 패널 PNG의 종횡비(폭 51.1 mm x 높이 18.6 mm)로 삽입 높이를 정한다.
PANEL_ASPECT = 22.0 / 50.2
PANEL_H = (COL_W - 2 * PAD) * PANEL_ASPECT
SLOT_H = 0.25 + LAB_H + PANEL_H + 0.2 + CAP_H
CARD_H = 2 * SLOT_H + 0.3

GOAL = "희소 라벨 조건에서의 광역 활동층 두께 예측 및 미관측 영역으로의 확장"
HEADS = ["기존 연구 한계점", "개선 방안", "기대 효과"]
HEAD_COL = [GREY_H, NAVY, MID]
CARDS = [
    [("1", "라벨 희소", "· 예측 대상 격자의 1.5%만 실측 · 74개 0.5° 블록에 집중"),
     ("2", "지역 간 전이 오차", "· 미관측 지역에서 데이터 기반 방법론 오차 급증")],
    [("3", "물리경험식 기반 라벨 증강", "· 물리식 정확도별 증강 효과 비교 · 증강 비율 반응 분석"),
     ("4", "물리 잔차 · 독립 관측 결합", "· 물리식 단독 대비 전이 오차 단계적 감소안 개발")],
    [("5", "광역 연속장 산출", "· 미관측 영역까지 활동층 두께 예측장 확장 가능"),
     ("6", "예측 신뢰 범위 제시", "· 보정 예측구간 폭과 전이 적용가능 영역 표기")],
]


def _txt(shape, text, size, color, bold=False, align=PP_ALIGN.LEFT, font=None):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.1)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.color.rgb = color
    f.name = font or FONT
    # 한글 대체 폰트 지정(Pretendard 미설치 환경 대비)
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag.split(':')[1]}")
        if el is None:
            from pptx.oxml.ns import qn
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", (font or FONT) if tag == "a:latin" else (font or FONT_FALLBACK))
    return shape


def _box(sl, x, y, w, h, fill, line, rounded=True):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Cm(x), Cm(y), Cm(w), Cm(h))
    shp.adjustments[0] = 0.04 if rounded else 0
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    shp.shadow.inherit = False
    return shp


def main():
    missing = [k for k in "123456" if not (PANEL / f"panel_{k}.png").exists()]
    if missing:
        raise SystemExit(f"패널 이미지 없음 {missing} — report_overview_figure.py 를 먼저 실행하라")

    prs = Presentation()
    prs.slide_width, prs.slide_height = Cm(SW), Cm(SH)
    sl = prs.slides.add_slide(prs.slide_layouts[6])       # 빈 레이아웃

    # 목표 띠
    y = M
    _box(sl, M, y, SW - 2 * M, GOAL_H, GOAL_BG, EDGE)
    _txt(_box(sl, M + 0.25, y + 0.25, 2.5, GOAL_H - 0.5, NAVY, None),
         "목표", 13, WHITE, bold=True, align=PP_ALIGN.CENTER)
    _txt(sl.shapes.add_textbox(Cm(M + 3.1), Cm(y), Cm(SW - 2 * M - 3.4), Cm(GOAL_H)),
         GOAL, 13, INK, bold=True, font="Pretendard ExtraBold")

    # 열
    y_head = y + GOAL_H + GAP
    for c in range(3):
        x0 = M + c * (COL_W + COL_GAP)
        _box(sl, x0, y_head, COL_W, HEAD_H, HEAD_COL[c], None)
        _txt(sl.shapes.add_textbox(Cm(x0), Cm(y_head), Cm(COL_W), Cm(HEAD_H)),
             HEADS[c], 14, WHITE, bold=True, align=PP_ALIGN.CENTER)

        y_card = y_head + HEAD_H + 0.3
        _box(sl, x0, y_card, COL_W, CARD_H, CARD_BG[c], CARD_EC[c])
        slot = SLOT_H
        for k, (key, lab, cap) in enumerate(CARDS[c]):
            top = y_card + k * slot + 0.25
            _txt(sl.shapes.add_textbox(Cm(x0 + PAD), Cm(top), Cm(COL_W - 2 * PAD), Cm(LAB_H)),
                 lab, 12.5, INK, bold=True)
            sl.shapes.add_picture(str(PANEL / f"panel_{key}.png"),
                                  Cm(x0 + PAD), Cm(top + LAB_H),
                                  width=Cm(COL_W - 2 * PAD))
            _txt(sl.shapes.add_textbox(Cm(x0 + PAD), Cm(top + LAB_H + PANEL_H + 0.15),
                                       Cm(COL_W - 2 * PAD), Cm(CAP_H)),
                 cap, 11, SUB, bold=True)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"[슬라이드] {SW:.2f} x {SH:.2f} cm (16:9) · 열 {COL_W:.2f} cm")
    print("saved", OUT.relative_to(ROOT))


if __name__ == "__main__":
    main()
