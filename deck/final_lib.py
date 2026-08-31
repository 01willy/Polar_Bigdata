# -*- coding: utf-8 -*-
"""final_lib.py — 본선 발표덱 디자인 시스템.

참고 스타일: IMAGE_SEUNGWONBAEK pseudo-MCMC 덱 (백색 배경 · 오렌지 악센트 · Arial 계열).
한글 대응: Pretendard (무료 OFL). 슬라이드 문법: 오렌지 볼드 타이틀 + 회색 부제 1줄
+ 그림 중심 본문 + 하단 한 줄 결론. 좌표 단위 inch, 캔버스 16:9 = 13.333 x 7.5.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

EMU_IN = 914400
SW, SH = 13.333, 7.5
ML, MR, MT = 0.52, 0.52, 0.30
CW = SW - ML - MR          # 12.293
CR = SW - MR
CONTENT_TOP = 1.18         # 헤더 아래 본문 상단
CONTENT_BOT = 6.80         # 결론 줄 위 본문 하한

# ---------- 팔레트 (참고덱 토큰) ----------
ACC   = RGBColor(0xEA, 0x85, 0x1B)   # 오렌지 악센트
ACC_DK= RGBColor(0xC0, 0x6A, 0x14)
INK   = RGBColor(0x11, 0x11, 0x11)
GRAY  = RGBColor(0x6B, 0x6B, 0x6B)
GRAY2 = RGBColor(0x4A, 0x4A, 0x4A)
MUTE  = RGBColor(0x8E, 0x8E, 0x8E)
HAIR  = RGBColor(0xDD, 0xDD, 0xDD)
PANEL = RGBColor(0xF7, 0xF7, 0xF5)
PANEL2= RGBColor(0xFB, 0xF1, 0xE4)   # 옅은 오렌지 (강조 패널)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
NAVY  = RGBColor(0x1F, 0x3A, 0x52)
STEEL = RGBColor(0x3D, 0x6B, 0x8E)

# ---------- 폰트 ----------
F_X = "Pretendard ExtraBold"
F_B = "Pretendard Bold"
F_S = "Pretendard SemiBold"
F_M = "Pretendard Medium"
F_R = "Pretendard"

CROP_DIR = "assets/final/crops"


def new_deck():
    prs = Presentation()
    prs.slide_width = Inches(SW)
    prs.slide_height = Inches(SH)
    return prs


def blank(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    return sl


def rect(sl, x, y, w, h, fill=PANEL, line=None, lw=0.75, shadow=False, round_=False,
         dash=None):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try:
            shp.adjustments[0] = 0.055
        except Exception:
            pass
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(lw)
        if dash:
            shp.line.dash_style = dash
    shp.shadow.inherit = False
    return shp


def hline(sl, x, y, w, color=HAIR, lw=0.9):
    ln = sl.shapes.add_connector(1, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(lw)
    ln.shadow.inherit = False
    return ln


def vline(sl, x, y, h, color=HAIR, lw=0.9):
    ln = sl.shapes.add_connector(1, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color; ln.line.width = Pt(lw)
    ln.shadow.inherit = False
    return ln


def text(sl, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.18, space_after=4, wrap=True):
    """paras: [run, ...] 또는 [[run, ...], ...]. run = dict(t,size,color,font,bold,spacing)."""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if paras and isinstance(paras[0], dict):
        paras = [paras]
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        for rd in para:
            r = p.add_run()
            r.text = rd["t"]
            r.font.size = Pt(rd.get("size", 12.5))
            r.font.color.rgb = rd.get("color", INK)
            _set_fonts(r, rd.get("font", F_M))
            if rd.get("bold"):
                r.font.bold = True
            if rd.get("underline"):
                r.font.underline = True
            if rd.get("sup"):
                r._r.get_or_add_rPr().set("baseline", "30000")
            if rd.get("spacing") is not None:
                try:
                    r.font._rPr.set("spc", str(int(rd["spacing"] * 100)))
                except Exception:
                    pass
    return tb


def _set_fonts(r, name):
    """라틴·한글(EA)·복합문자(CS) 폰트를 모두 지정한다. 미지정 시 한글이
    테마 기본 폰트로 대체되어 문장부호 간격이 벌어진다."""
    r.font.name = name
    rPr = r._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def pagenum(sl, n):
    text(sl, SW - 0.95, 0.22, 0.6, 0.3, [{"t": str(n), "size": 12, "color": GRAY, "font": F_M}],
         align=PP_ALIGN.RIGHT)


def header(sl, title, subtitle=None, n=None):
    """오렌지 볼드 타이틀 + 회색 부제. 본문 상단(CONTENT_TOP) 반환."""
    text(sl, ML, 0.24, CW - 0.7, 0.55,
         [{"t": title, "size": 26, "color": ACC, "font": F_X}])
    if subtitle:
        text(sl, ML + 0.04, 0.80, CW - 0.8, 0.36,
             [{"t": subtitle, "size": 14, "color": GRAY, "font": F_M}])
    if n is not None:
        pagenum(sl, n)
    return CONTENT_TOP


def takeaway(sl, msg, y=7.02, accent=False):
    """하단 한 줄 결론: 중앙 정렬(참고덱 문법). accent=True면 잉크 볼드 확대."""
    hline(sl, ML, y - 0.10, CW, color=HAIR, lw=0.9)
    text(sl, ML + 0.02, y, CW, 0.4,
         [{"t": msg, "size": 15 if accent else 14.5, "color": INK,
           "font": F_X if accent else F_S}],
         align=PP_ALIGN.CENTER)


def section_label(sl, x, y, t, color=ACC, size=14):
    text(sl, x, y, 6.0, 0.32, [{"t": t, "size": size, "color": color, "font": F_X}])


def bullets(sl, x, y, w, items, size=12.5, gap=6, color=INK, dot_color=ACC,
            line_spacing=1.22, sub_size=11, sub_color=GRAY):
    """items: [(본문, 보조설명|None), ...]. 오렌지 불릿 점."""
    paras = []
    for main, sub in items:
        paras.append([{"t": "●  ", "size": size - 4.5, "color": dot_color, "font": F_S},
                      {"t": main, "size": size, "color": color, "font": F_M}])
        if sub:
            paras.append([{"t": "     " + sub, "size": sub_size, "color": sub_color, "font": F_M}])
    text(sl, x, y, w, 4.5, paras, line_spacing=line_spacing, space_after=gap)


def _img_size(path):
    with Image.open(path) as im:
        return im.size


def crop_asset(src, box_px, name):
    """src를 픽셀 박스(l,t,r,b)로 잘라 크롭 자산으로 저장하고 경로 반환."""
    os.makedirs(CROP_DIR, exist_ok=True)
    out = os.path.join(CROP_DIR, name)
    if not os.path.exists(out):
        with Image.open(src) as im:
            im.crop(box_px).save(out)
    return out


def image(sl, path, x, y, w=None, h=None, frame=False):
    """가로폭 또는 세로 기준 배치. 실제 (w, h) 반환."""
    iw, ih = _img_size(path)
    ar = iw / ih
    if w is not None:
        h = w / ar
    elif h is not None:
        w = h * ar
    pic = sl.shapes.add_picture(path, Inches(x), Inches(y), Inches(w), Inches(h))
    pic.shadow.inherit = False
    if frame:
        pic.line.color.rgb = HAIR; pic.line.width = Pt(0.75)
    return w, h


def caption(sl, x, y, w, t, align=PP_ALIGN.LEFT):
    text(sl, x, y, w, 0.3, [{"t": t, "size": 10.5, "color": MUTE, "font": F_M}], align=align)


def stat_card(sl, x, y, w, h, tag, title, value, note, tag_color=ACC):
    """참고덱 결과 카드 스타일: 옅은 패널 + 태그 + 수치 강조."""
    rect(sl, x, y, w, h, fill=PANEL, line=HAIR, lw=0.75, round_=True)
    text(sl, x + 0.18, y + 0.14, w - 0.36, 0.3,
         [{"t": tag + "  ", "size": 12.5, "color": tag_color, "font": F_X},
          {"t": title, "size": 12.5, "color": INK, "font": F_S}])
    text(sl, x + 0.18, y + 0.52, w - 0.36, 0.42,
         [{"t": value, "size": 19, "color": INK, "font": F_X}])
    text(sl, x + 0.18, y + 1.02, w - 0.36, 0.3,
         [{"t": note, "size": 10.5, "color": GRAY, "font": F_M}])


def mini_table(sl, x, y, w, rows, col_w=None, size=12, header_row=True,
               row_h=0.34, align_cols=None, header_color=GRAY2, pad=0.06):
    """booktabs 풍 소형 표: 상단·헤더 아래·하단 규칙선."""
    n_col = len(rows[0])
    if col_w is None:
        col_w = [w / n_col] * n_col
    if align_cols is None:
        align_cols = [PP_ALIGN.LEFT] + [PP_ALIGN.RIGHT] * (n_col - 1)
    yy = y
    hline(sl, x, yy, w, color=GRAY2, lw=1.1)
    yy += pad
    for ri, row in enumerate(rows):
        xx = x
        is_head = header_row and ri == 0
        for ci, cell in enumerate(row):
            if isinstance(cell, dict):
                rd = dict(cell)
            else:
                rd = {"t": str(cell)}
            rd.setdefault("size", size - 0.5 if is_head else size)
            rd.setdefault("color", header_color if is_head else INK)
            rd.setdefault("font", F_S if is_head else F_M)
            text(sl, xx + 0.02, yy + 0.02, col_w[ci] - 0.06, row_h, [rd],
                 align=align_cols[ci])
            xx += col_w[ci]
        yy += row_h
        if is_head:
            hline(sl, x, yy, w, color=HAIR, lw=0.8)
            yy += pad
    hline(sl, x, yy + 0.02, w, color=GRAY2, lw=1.1)
    return yy + 0.02
