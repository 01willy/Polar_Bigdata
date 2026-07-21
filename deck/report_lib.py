# -*- coding: utf-8 -*-
"""
report_lib.py — Polar_Bigdata 발표덱 v2 디자인 시스템 (에디토리얼/학술 보고서 톤).

설계 의도(v1의 "AI틱 라운드카드·테크그라디언트" 결별):
- 종이 배경(warm off-white) · 세리프 제목(Noto Serif CJK KR) + 산세리프 본문(Pretendard).
- 라운드/그림자/채운카드 금지. booktabs 표(가로 3선), 번호 캡션(그림 N.), 얇은 규칙선.
- 냉색 최소 팔레트(딥틸 1 accent + 슬레이트). 수식 렌더링 지원.
좌표 단위 inch, 캔버스 16:9 = 13.333 x 7.5.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

EMU_IN = 914400
SW, SH = 13.333, 7.5
# 여백(학술 페이지: 넉넉한 좌/우, 러닝헤더·푸터)
ML, MR, MT, MB = 0.85, 0.85, 0.62, 0.5
CW = SW - ML - MR
CR = SW - MR

# ---------- 팔레트 (warm-neutral paper + cool accent) ----------
PAPER   = RGBColor(0xFB, 0xFA, 0xF6)   # 종이 배경
INK     = RGBColor(0x1B, 0x1E, 0x23)   # 본문 잉크
INK2    = RGBColor(0x3A, 0x40, 0x47)   # 부제
SLATE   = RGBColor(0x5C, 0x66, 0x6E)   # 보조
MUTE    = RGBColor(0x8A, 0x90, 0x96)   # 캡션/약함
TEAL    = RGBColor(0x0E, 0x5A, 0x61)   # 주 accent(딥틸, 냉색)
TEAL2   = RGBColor(0x2C, 0x7D, 0x83)   # 밝은 틸
NAVY    = RGBColor(0x1F, 0x3A, 0x52)   # 짙은 청(데이터)
GOLD    = RGBColor(0x9A, 0x6D, 0x1E)   # 강조 소량(문자에만, 데이터장 아님)
RULE    = RGBColor(0xD9, 0xD6, 0xCC)   # 헤어라인
RULE2   = RGBColor(0xBE, 0xBA, 0xAE)   # 진한 규칙선
FILL    = RGBColor(0xF2, 0xF0, 0xE9)   # 아주 옅은 패널(드묾)
INK_DK  = RGBColor(0x14, 0x28, 0x30)   # 표지 딥
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ---------- 폰트 (Pretendard, ~/.fonts 설치 확인 · LibreOffice fontconfig 렌더) ----------
SERIF   = "Pretendard ExtraBold"  # 제목(굵은 산세리프, EMP 덱 톤)
SANS    = "Pretendard SemiBold"   # 본문
SANS_M  = "Pretendard Medium"
SANS_S  = "Pretendard SemiBold"
SANS_L  = "Pretendard"


def new_deck():
    p = Presentation(); p.slide_width = Emu(int(SW*EMU_IN)); p.slide_height = Emu(int(SH*EMU_IN)); return p

def slide(prs, bg=PAPER):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Emu(int(SW*EMU_IN)), Emu(int(SH*EMU_IN)))
    r.fill.solid(); r.fill.fore_color.rgb = bg; r.line.fill.background(); r.shadow.inherit = False
    return s

def _ea(rPr, font):
    for tag in ('a:ea', 'a:cs'):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {}); rPr.append(el)
        el.set('typeface', font)

def _run(p, t, size, color, font=SANS, bold=False, italic=False, spc=None):
    if t.endswith('.') and not t.endswith('..'):
        t = t[:-1]
    r = p.add_run(); r.text = t
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = font
    r.font.bold = bold; r.font.italic = italic
    rPr = r._r.get_or_add_rPr(); _ea(rPr, font)
    if spc is not None:
        rPr.set('spc', str(int(spc*100)))
    return r

def text(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.1, sa=2.0, wrap=True):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = wrap; tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, 0)
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls; p.space_after = Pt(sa); p.space_before = Pt(0)
        if isinstance(para, dict): para = [para]
        for rd in para:
            _run(p, rd.get('t',''), rd.get('size',13), rd.get('color',INK), rd.get('font',SANS),
                 rd.get('bold',False), rd.get('italic',False), rd.get('spc',None))
    return tb

def rule(s, x, y, w, color=RULE, wt=1.0):
    ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x+w), Inches(y))
    ln.line.color.rgb = color; ln.line.width = Pt(wt); ln.shadow.inherit = False; return ln

def vrule(s, x, y, h, color=RULE, wt=1.0):
    ln = s.shapes.add_connector(2, Inches(x), Inches(y), Inches(x), Inches(y+h))
    ln.line.color.rgb = color; ln.line.width = Pt(wt); ln.shadow.inherit = False; return ln

def rect(s, x, y, w, h, fill=None, line=None, wt=1.0):
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb = fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb = line; sp.line.width = Pt(wt)
    sp.shadow.inherit = False; return sp

def _aspect(path):
    with Image.open(path) as im:
        return im.size[0]/im.size[1]

def image(s, path, x, y, w, h, align='center', valign='middle'):
    a = _aspect(path); ba = w/h
    if a >= ba: dw, dh = w, w/a
    else: dh, dw = h, h*a
    dx = x + (w-dw)/2 if align=='center' else (x if align=='left' else x+w-dw)
    dy = y + (h-dh)/2 if valign=='middle' else (y if valign=='top' else y+h-dh)
    s.shapes.add_picture(path, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    return (dx, dy, dw, dh)

# ---------- 러닝 헤더/푸터 (저널 페이지 느낌) ----------
def chrome(s, idx, total, section, running="영구동토 활동층 두께 예측 · 3D 지중온도장, 관측기반 GeoAI"):
    # 상단 러닝헤더
    text(s, ML, 0.32, CW*0.7, 0.24, [[{'t': running, 'size': 8, 'color': MUTE, 'font': SANS_M, 'spc': 0.3}]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, CR-3.0, 0.32, 3.0, 0.24, [[{'t': section, 'size': 8, 'color': MUTE, 'font': SANS_M, 'spc': 0.3}]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
    rule(s, ML, 0.56, CW, color=RULE, wt=0.75)
    # 하단 푸터
    yb = SH-0.4
    rule(s, ML, yb-0.04, CW, color=RULE, wt=0.75)
    text(s, ML, yb, CW*0.6, 0.24, [[{'t': "Polar_Bigdata · 진행 보고 (2026-07)", 'size': 8, 'color': MUTE, 'font': SANS_M}]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, CR-1.2, yb, 1.2, 0.24, [[{'t': f"{idx} / {total}", 'size': 8.5, 'color': SLATE, 'font': SANS_M}]],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def title_block(s, no, kicker, title, sub=None):
    """섹션번호 · 키커 · 세리프 제목 · (부제). 콘텐츠 시작 y 반환."""
    y = 0.82
    text(s, ML, y, CW, 0.24,
         [[{'t': f"{no}", 'size': 10.5, 'color': TEAL, 'font': SANS_S, 'spc': 0.8},
           {'t': "   " + kicker, 'size': 10.5, 'color': SLATE, 'font': SANS_M, 'spc': 0.8}]],
         anchor=MSO_ANCHOR.MIDDLE)
    text(s, ML, y+0.3, CW, 0.62, [[{'t': title, 'size': 25, 'color': INK, 'font': SERIF, 'bold': True}]],
         anchor=MSO_ANCHOR.MIDDLE, ls=1.0)
    yy = y+0.98
    if sub:
        text(s, ML, yy, CW, 0.34, [[{'t': sub, 'size': 12.5, 'color': INK2, 'font': SANS}]], ls=1.15)
        yy += 0.4
    rule(s, ML, yy, CW, color=RULE2, wt=1.3)
    return yy+0.18

def caption(s, x, y, w, num, txt, src=None, align=PP_ALIGN.LEFT):
    paras = [[{'t': f"그림 {num}. ", 'size': 9.5, 'color': TEAL, 'font': SANS_S},
              {'t': txt, 'size': 9.5, 'color': SLATE, 'font': SANS_M}]]
    if src:
        paras.append([{'t': src, 'size': 8, 'color': MUTE, 'font': SANS}])
    text(s, x, y, w, 0.5, paras, align=align, ls=1.12, sa=1.0)

def tcaption(s, x, y, w, num, txt, src=None):
    paras = [[{'t': f"표 {num}. ", 'size': 9.5, 'color': TEAL, 'font': SANS_S},
              {'t': txt, 'size': 9.5, 'color': SLATE, 'font': SANS_M}]]
    if src: paras.append([{'t': src, 'size': 8, 'color': MUTE, 'font': SANS}])
    text(s, x, y, w, 0.4, paras, ls=1.1, sa=1.0)

# ---------- booktabs 표 (가로 3선, 세로선 없음, zebra 없음) ----------
def booktable(s, x, y, w, rows, col_w, header=True, size=11, row_h=0.34, head_h=0.36,
              align=None, hi_rows=None, hi_color=None):
    n, c = len(rows), len(rows[0])
    hi_rows = hi_rows or {}
    total_h = (head_h if header else 0) + (n-(1 if header else 0))*row_h
    gt = s.shapes.add_table(n, c, Inches(x), Inches(y), Inches(w), Inches(total_h)); tbl = gt.table
    tbl._tbl.tblPr.set('firstRow','0'); tbl._tbl.tblPr.set('bandRow','0')
    for j, cw in enumerate(col_w): tbl.columns[j].width = Inches(cw)
    for i in range(n): tbl.rows[i].height = Inches(head_h if (header and i==0) else row_h)
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = PAPER
            for m in ('margin_left','margin_right','margin_top','margin_bottom'):
                setattr(cell, m, Inches(0.03 if m in ('margin_top','margin_bottom') else 0.07))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]
            a = (align[j] if align else (PP_ALIGN.LEFT if j==0 else PP_ALIGN.CENTER))
            p.alignment = a
            ishead = header and i==0
            col = INK if not ishead else INK
            fnt = SANS_S if (ishead or j==0) else SANS
            if i in hi_rows and not ishead:
                cell.fill.fore_color.rgb = FILL
                col = hi_color or TEAL; fnt = SANS_S
            _run(p, str(val), size, col, fnt, bold=False)
    # booktabs 규칙선: 상단 굵게, 헤더 아래, 하단 굵게
    rule(s, x, y-0.005, w, color=INK, wt=1.6)
    if header:
        rule(s, x, y+head_h, w, color=RULE2, wt=0.9)
    rule(s, x, y+total_h+0.005, w, color=INK, wt=1.6)
    return total_h

# ---------- 수치 요약 (라벨 위 · 숫자 아래, 장식 규칙선 없음) ----------
def fom(s, x, y, w, value, unit, label, sub=None, color=INK, accent=TEAL):
    text(s, x, y, w, 0.24, [[{'t': label, 'size': 9.5, 'color': MUTE, 'font': SANS_M, 'spc': 0.2}]])
    text(s, x, y+0.24, w, 0.52,
         [[{'t': value, 'size': 26, 'color': color, 'font': SERIF, 'bold': True},
           {'t': (" "+unit if unit else ""), 'size': 12.5, 'color': SLATE, 'font': SANS_M}]],
         anchor=MSO_ANCHOR.MIDDLE)
    if sub:
        text(s, x, y+0.82, w, 0.3, [[{'t': sub, 'size': 9, 'color': SLATE, 'font': SANS}]], ls=1.12)

# ---------- 라벨형 설명(장식 바 없이 라벨 강조 + 본문) ----------
def finding(s, x, y, w, label, body, accent=TEAL):
    text(s, x, y, w, 0.7,
         [[{'t': label + "  ", 'size': 11.5, 'color': accent, 'font': SANS_S},
           {'t': body, 'size': 11.5, 'color': INK2, 'font': SANS}]], ls=1.24, sa=1.0)

def bullets(s, x, y, w, items, size=12.5, gap=8, color=INK2, head=INK, lead=TEAL):
    """(head, body) 항목은 head를 accent 라벨로 run-in. 장식 대시 없음(보고서 정의형)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(3)); tf = tb.text_frame
    tf.word_wrap = True
    for m in ('margin_left','margin_right','margin_top','margin_bottom'): setattr(tf, m, 0)
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.line_spacing = 1.2; p.space_after = Pt(gap); p.space_before = Pt(0)
        if isinstance(it, (list, tuple)):
            _run(p, it[0], size, lead, SANS_S)
            if it[1]: _run(p, "   " + it[1], size, color, SANS)
        else:
            _run(p, "·  ", size, lead, SANS_S)
            _run(p, it, size, color, SANS)
    return tb

# ---------- 수식 렌더링 (matplotlib mathtext → 투명 PNG) ----------
_EQ_CACHE = "assets/eq"  # build는 deck/에서 실행 → deck/assets/eq/
def equation(latex, name, fontsize=22, color="#1B1E23", dpi=300):
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.pyplot as plt
    os.makedirs(_EQ_CACHE, exist_ok=True)
    path = f"{_EQ_CACHE}/{name}.png"
    fig = plt.figure(figsize=(0.01,0.01))
    fig.text(0, 0, f"${latex}$", fontsize=fontsize, color=color)
    fig.savefig(path, dpi=dpi, transparent=True, bbox_inches="tight", pad_inches=0.14)
    plt.close(fig)
    return path
