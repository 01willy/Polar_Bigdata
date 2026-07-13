# -*- coding: utf-8 -*-
"""
polar_slides.py — Polar_Bigdata 발표덱용 재사용 레이아웃/테마 라이브러리 (python-pptx).

디자인 규약 근거: design/brand_tokens.json, design/layout_rules.md.
- 냉색 표준 팔레트(붉은 계열 금지), Pretendard 타이포, 격자 정렬, 보고서 톤.
- 슬라이드 크롬(제목/키커/규칙선/푸터)은 Pretendard, 그림은 원본(matplotlib NanumGothic) 유지.

모든 좌표 단위는 inch. 캔버스 16:9 = 13.333 x 7.5 in.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn
from PIL import Image

# ----------------------------------------------------------------------------
# 캔버스
# ----------------------------------------------------------------------------
EMU_PER_IN = 914400
SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

# 여백 그리드
ML = 0.72          # left margin
MR = 0.72          # right margin
MT = 0.62          # top margin
MB = 0.46          # bottom margin
CONTENT_W = SLIDE_W_IN - ML - MR
CONTENT_L = ML
CONTENT_R = SLIDE_W_IN - MR

# ----------------------------------------------------------------------------
# 색 (design/brand_tokens.json categorical + 파생 중성색)
# ----------------------------------------------------------------------------
INK      = RGBColor(0x16, 0x22, 0x2E)   # 본문 잉크(짙은 네이비블랙)
NAVY     = RGBColor(0x1F, 0x4E, 0x79)   # gbm / primary
TEAL     = RGBColor(0x0B, 0x72, 0x85)   # highlight / accent
TEAL_LT  = RGBColor(0x2E, 0x8B, 0x9E)   # dl
PURPLE   = RGBColor(0x6A, 0x51, 0xA3)   # ensemble
WARN     = RGBColor(0xB5, 0x65, 0x1D)   # 주의(정정/caveat)
SLATE    = RGBColor(0x54, 0x60, 0x6B)   # 보조 텍스트
GRAY     = RGBColor(0x8A, 0x8F, 0x98)   # baseline / muted
HAIR     = RGBColor(0xD5, 0xDB, 0xE1)   # 헤어라인 규칙선
HAIR_LT  = RGBColor(0xE7, 0xEB, 0xEF)
PANEL    = RGBColor(0xF4, 0xF6, 0xF8)   # 옅은 패널 배경
PANEL2   = RGBColor(0xEC, 0xF1, 0xF4)   # 옅은 청색 패널
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
INK_ON   = RGBColor(0xF2, 0xF6, 0xF9)   # 어두운 배경 위 밝은 텍스트
NAVY_DK  = RGBColor(0x14, 0x33, 0x50)   # 표지/디바이더 딥네이비
NAVY_DK2 = RGBColor(0x0E, 0x26, 0x3D)

# ----------------------------------------------------------------------------
# 폰트 (Pretendard: 모던·전문·보고서톤)
# ----------------------------------------------------------------------------
F_SEMI  = "Pretendard SemiBold"   # 제목/헤더/강조 라벨
F_MED   = "Pretendard Medium"     # 서브·키커
F_REG   = "Pretendard"            # 본문
F_LIGHT = "Pretendard Light"      # 캡션/보조


# ----------------------------------------------------------------------------
# 저수준 헬퍼
# ----------------------------------------------------------------------------
def new_deck():
    prs = Presentation()
    prs.slide_width = Emu(int(SLIDE_W_IN * EMU_PER_IN))
    prs.slide_height = Emu(int(SLIDE_H_IN * EMU_PER_IN))
    return prs


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _no_line(shape):
    shape.line.fill.background()


def _solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    _no_line(shape)


def rect(slide, x, y, w, h, color=None, line=None, line_w=0.75, shadow=False):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if color is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if not shadow:
        pass
    return sp


def hline(slide, x, y, w, color=HAIR, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def vline(slide, x, y, h, color=HAIR, weight=1.0):
    ln = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x), Inches(y + h))
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def _set_run(r, text, size, color, font=F_REG, bold=False, italic=False, spacing=None):
    # 슬라이드 관용: 종결 마침표 제거. LibreOffice가 한글↔라틴 경계(예 "다.")에
    # 삽입하는 CJK–서양 자동간격 아티팩트를 함께 제거한다(폰트로는 못 없앰).
    if text.endswith('.') and not text.endswith('..'):
        text = text[:-1]
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.name = font
    r.font.bold = bold
    r.font.italic = italic
    # East-Asian 폰트도 동일 지정(한글 정상 적용)
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', font)
    cs = rPr.find(qn('a:cs'))
    if cs is None:
        cs = rPr.makeelement(qn('a:cs'), {})
        rPr.append(cs)
    cs.set('typeface', font)
    if spacing is not None:
        rPr.set('spc', str(int(spacing * 100)))  # pt*100


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.06, wrap=True, space_after=2.0):
    """runs: list of paragraphs; each paragraph = list of run-dicts
       run-dict: {t, size, color, font, bold, italic, spacing}
       또는 문단이 단일 dict면 단일 run 문단."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        if isinstance(para, dict):
            para = [para]
        for j, rd in enumerate(para):
            r = p.add_run()
            _set_run(r, rd.get('t', ''), rd.get('size', 14), rd.get('color', INK),
                     rd.get('font', F_REG), rd.get('bold', False), rd.get('italic', False),
                     rd.get('spacing', None))
    return tb


def _aspect(path):
    with Image.open(path) as im:
        w, h = im.size
    return w / h


def image_fit(slide, path, x, y, w, h, align='center', valign='middle'):
    """박스(x,y,w,h) 안에 종횡비 유지하며 최대로 배치. 실제 배치 박스 반환."""
    a = _aspect(path)
    box_a = w / h
    if a >= box_a:            # 이미지가 더 넓음 → 폭 맞춤
        dw = w
        dh = w / a
    else:                     # 이미지가 더 높음 → 높이 맞춤
        dh = h
        dw = h * a
    if align == 'center':
        dx = x + (w - dw) / 2
    elif align == 'left':
        dx = x
    else:
        dx = x + (w - dw)
    if valign == 'middle':
        dy = y + (h - dh) / 2
    elif valign == 'top':
        dy = y
    else:
        dy = y + (h - dh)
    slide.shapes.add_picture(path, Inches(dx), Inches(dy), Inches(dw), Inches(dh))
    return (dx, dy, dw, dh)


# ----------------------------------------------------------------------------
# 고수준 레이아웃 프리미티브
# ----------------------------------------------------------------------------
def footer(slide, idx, total, tag="극지 빅데이터 · 영구동토 ALT/3D 열구조"):
    y = SLIDE_H_IN - 0.36
    hline(slide, ML, y - 0.06, CONTENT_W, color=HAIR_LT, weight=0.75)
    text(slide, ML, y, CONTENT_W * 0.7, 0.28,
         [{'t': tag, 'size': 8, 'color': GRAY, 'font': F_MED}],
         anchor=MSO_ANCHOR.MIDDLE)
    text(slide, CONTENT_R - 1.6, y, 1.6, 0.28,
         [{'t': f"{idx:02d} / {total:02d}", 'size': 8, 'color': GRAY, 'font': F_MED}],
         align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def header(slide, kicker, title, accent=TEAL):
    """상단 제목 밴드. 콘텐츠 시작 y를 반환."""
    # kicker (섹션 라벨)
    text(slide, ML, MT, CONTENT_W, 0.26,
         [[{'t': "▍ ", 'size': 11, 'color': accent, 'font': F_SEMI},
           {'t': kicker, 'size': 11, 'color': accent, 'font': F_SEMI, 'spacing': 0.6}]],
         anchor=MSO_ANCHOR.MIDDLE)
    # title
    text(slide, ML, MT + 0.28, CONTENT_W, 0.62,
         [{'t': title, 'size': 24, 'color': INK, 'font': F_SEMI}],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.0)
    y_rule = MT + 0.98
    hline(slide, ML, y_rule, CONTENT_W, color=HAIR, weight=1.2)
    return y_rule + 0.16


def title_slide(prs, title, subtitle, meta_lines, hero_path=None, phase_tag=None):
    slide = blank(prs)
    # 딥네이비 배경
    rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, color=NAVY_DK)
    # 우측 히어로 이미지(있으면) — 밴드
    if hero_path:
        band_x = 8.05
        rect(slide, band_x, 0, SLIDE_W_IN - band_x, SLIDE_H_IN, color=NAVY_DK2)
        image_fit(slide, hero_path, band_x + 0.15, 0.5, SLIDE_W_IN - band_x - 0.35, SLIDE_H_IN - 1.0,
                  align='center', valign='middle')
        vline(slide, band_x, 0, SLIDE_H_IN, color=RGBColor(0x27, 0x48, 0x66), weight=1.5)
    # 상단 얇은 액센트
    rect(slide, ML, 1.35, 0.62, 0.09, color=TEAL)
    if phase_tag:
        text(slide, ML, 1.02, 7.0, 0.3,
             [{'t': phase_tag, 'size': 12, 'color': TEAL_LT, 'font': F_SEMI, 'spacing': 1.2}])
    tw = 7.0 if hero_path else 11.0
    text(slide, ML, 1.62, tw, 2.2,
         [{'t': title, 'size': 38, 'color': WHITE, 'font': F_SEMI}], line_spacing=1.04)
    text(slide, ML, 3.7, tw, 1.2,
         [{'t': subtitle, 'size': 16, 'color': RGBColor(0xC3, 0xD1, 0xDD), 'font': F_REG}],
         line_spacing=1.22)
    # 메타(하단)
    hline(slide, ML, 6.35, tw - 0.2, color=RGBColor(0x2A, 0x4B, 0x69), weight=1.0)
    metas = []
    for ln in meta_lines:
        metas.append([{'t': ln, 'size': 11.5, 'color': RGBColor(0x9F, 0xB4, 0xC6), 'font': F_MED}])
    text(slide, ML, 6.5, tw - 0.2, 0.8, metas, line_spacing=1.25, space_after=1.5)
    return slide


def section_divider(prs, no, title, subtitle=None, idx=None, total=None):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W_IN, SLIDE_H_IN, color=NAVY_DK)
    rect(slide, 0, 0, 0.28, SLIDE_H_IN, color=TEAL)
    text(slide, ML + 0.1, 2.55, 3.0, 1.6,
         [{'t': no, 'size': 82, 'color': RGBColor(0x2E, 0x8B, 0x9E), 'font': F_SEMI}])
    text(slide, ML + 2.35, 2.9, 9.4, 1.2,
         [{'t': title, 'size': 34, 'color': WHITE, 'font': F_SEMI}], line_spacing=1.02)
    if subtitle:
        text(slide, ML + 2.4, 3.95, 9.3, 1.0,
             [{'t': subtitle, 'size': 14, 'color': RGBColor(0xAD, 0xC1, 0xD1), 'font': F_REG}],
             line_spacing=1.25)
    hline(slide, ML + 2.4, 3.78, 6.0, color=RGBColor(0x2A, 0x4B, 0x69), weight=1.0)
    return slide


def bullets(slide, x, y, w, h, items, size=14.5, gap=7.0, lead_color=TEAL,
            title_color=INK, body_color=SLATE):
    """items: list of (head, body) 또는 문자열. 편집형 대시 불릿."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.12
        p.space_after = Pt(gap)
        p.space_before = Pt(0)
        r0 = p.add_run()
        _set_run(r0, "—  ", size, lead_color, F_SEMI)
        if isinstance(it, (list, tuple)):
            head, body = it
            rh = p.add_run()
            _set_run(rh, head, size, title_color, F_SEMI)
            if body:
                rb = p.add_run()
                _set_run(rb, "  " + body, size, body_color, F_REG)
        else:
            rr = p.add_run()
            _set_run(rr, it, size, body_color, F_REG)
    return tb


def caption(slide, x, y, w, txt, source=None, size=9.5, align=PP_ALIGN.LEFT):
    runs = [{'t': txt, 'size': size, 'color': SLATE, 'font': F_MED}]
    paras = [runs]
    if source:
        paras.append([{'t': "근거 · " + source, 'size': size - 1,
                       'color': RGBColor(0x6a, 0x74, 0x80), 'font': F_REG}])
    text(slide, x, y, w, 0.6, paras, align=align, line_spacing=1.12, space_after=1.0)


def panel(slide, x, y, w, h, fill=PANEL, line=HAIR, line_w=1.0, radius=True):
    shape_t = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(shape_t, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if radius:
        try:
            sp.adjustments[0] = 0.045
        except Exception:
            pass
    return sp


def kpi(slide, x, y, w, h, value, unit, label, accent=TEAL, sub=None):
    """단일 KPI 카드."""
    panel(slide, x, y, w, h, fill=WHITE, line=HAIR, line_w=1.0)
    rect(slide, x, y + 0.14, 0.07, h - 0.28, color=accent)
    text(slide, x + 0.24, y + 0.16, w - 0.34, 0.7,
         [[{'t': value, 'size': 30, 'color': INK, 'font': F_SEMI},
           {'t': ("  " + unit if unit else ""), 'size': 13, 'color': SLATE, 'font': F_MED}]],
         anchor=MSO_ANCHOR.MIDDLE)
    yy = y + 0.86
    text(slide, x + 0.24, yy, w - 0.34, 0.34,
         [{'t': label, 'size': 11.5, 'color': NAVY, 'font': F_SEMI}], line_spacing=1.02)
    if sub:
        text(slide, x + 0.24, yy + 0.3, w - 0.34, 0.5,
             [{'t': sub, 'size': 9.5, 'color': SLATE, 'font': F_REG}], line_spacing=1.08)


def styled_table(slide, x, y, w, rows, col_w=None, header=True,
                 head_fill=NAVY, head_txt=WHITE, size=11.5, row_h=0.34,
                 head_h=0.4, zebra=True, align_first_left=True):
    """rows: list of list(str). 첫 행이 헤더(header=True)."""
    nrows = len(rows)
    ncols = len(rows[0])
    total_h = head_h + (nrows - 1) * row_h if header else nrows * row_h
    gtbl = slide.shapes.add_table(nrows, ncols, Inches(x), Inches(y), Inches(w), Inches(total_h))
    tbl = gtbl.table
    # 컬럼 폭
    if col_w:
        for c, cw in enumerate(col_w):
            tbl.columns[c].width = Inches(cw)
    # 행 높이
    for rIdx in range(nrows):
        tbl.rows[rIdx].height = Inches(head_h if (header and rIdx == 0) else row_h)
    # 스타일 제거(기본 테마 밴딩 off)
    tblPr = tbl._tbl.tblPr
    tblPr.set('firstRow', '0')
    tblPr.set('bandRow', '0')
    for rIdx, row in enumerate(rows):
        for cIdx, val in enumerate(row):
            cell = tbl.cell(rIdx, cIdx)
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.08)
            cell.margin_top = Inches(0.02)
            cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            is_head = header and rIdx == 0
            if is_head:
                cell.fill.solid(); cell.fill.fore_color.rgb = head_fill
            elif zebra and (rIdx % 2 == 0):
                cell.fill.solid(); cell.fill.fore_color.rgb = PANEL
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = WHITE
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if (cIdx == 0 and align_first_left) else PP_ALIGN.CENTER
            r = p.add_run()
            col = head_txt if is_head else INK
            fnt = F_SEMI if is_head else (F_SEMI if cIdx == 0 else F_REG)
            _set_run(r, val, size, col, fnt)
    # 헤어라인 테두리
    return gtbl


def chip(slide, x, y, w, h, label, fill=PANEL2, txt=NAVY, line=None, size=11, font=F_SEMI):
    sp = panel(slide, x, y, w, h, fill=fill, line=line, line_w=1.0)
    text(slide, x + 0.06, y, w - 0.12, h,
         [{'t': label, 'size': size, 'color': txt, 'font': font}],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.02)
    return sp


def arrow(slide, x, y, w, h=0.3, color=GRAY):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    _solid(sp, color)
    sp.shadow.inherit = False
    try:
        sp.adjustments[0] = 0.5
        sp.adjustments[1] = 0.55
    except Exception:
        pass
    return sp
